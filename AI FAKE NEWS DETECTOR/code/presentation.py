import os
import tkinter as tk
from tkinter import ttk, messagebox
from pptx import Presentation

class PowerPointNotesManager:
    """Handles adding speaker notes to PowerPoint slides"""
    def __init__(self):
        self.input_ppt = os.path.abspath("fake_news_detector/Presentation - AI Fake News Detector -1.pptx")
        self.output_ppt = os.path.abspath("AI_Fake_News_Detector_with_Notes.pptx")
        # Complete list of all 14 slide notes
        self.slide_notes = [
            # Slide 1: Title Slide
            "Hello everyone,This is our mini project – AI Fake News Detector.It’s a simple Python tool to identify fake news headlines.I’m Ravichandra, and along with my team — Sravya, Kalpana, and Latha Sri — we’re excited to show you our work.Let’s get started!",
                        
            # Slide 2: Problem Statement
            "Fake news spreads fast and affects public trust.It often goes viral through misinformation Our goal was to build a tool that tackles this problem.We used a simple method: keyword scoring, confidence percentage, and a user-friendly GUI — no internet or ML needed.",           
            # Slide 3: Team Introduction
            "Fake news causes real damage — it spreads misinformation, Our goal was to build a tool that tackles this problem.We used a simple method: keyword scoring, confidence percentage, and a user-friendly GUI — no internet or ML needed.",
            
            # Slide 4: Project Overview
            "Our tool classifies news headlines as Fake, Real, or Uncertain,and also shows a confidence percentage for each result.This helps users quickly judge if the news they see online is trustworthy.",
            
            # Slide 5: Why This Matters
           "We used Python 3 with Tkinter to build the GUI.Our dataset uses custom-weighted keywords to detect whether a headline is fake or real.The setup is light, simple, and doesn’t need any external libraries or machine learning models.",
            
            # Slide 6: Technology Stack
            "This project was a true team effort.Each member played a key role—from dataset research to GUI development and testing.Together, we planned, coded, reviewed, and refined every step to make this project functional and user-friendly.",
            
            # Slide 7: Keyword Methodology
            "We created weighted dictionaries of 50+ fake news indicators (like 'shocking', 'miracle') "
            "and 40+ real news markers (like 'official', 'verified').",
            
            # Slide 8: Scoring System
            "Each keyword has a weight from 1-3. The system calculates a confidence score "
            "based on the ratio of fake to real keywords detected.",
            
            # Slide 9: User Interface
            "The GUI features color-coded results: Red for fake, green for real, "
            "and yellow for uncertain classifications.",
            
            # Slide 10: Testing Process
            "We validated against 200 sample headlines with 85% accuracy. "
            "False positives occur mainly with satirical content.",
            
            # Slide 11: Integration
            "All components work together seamlessly: input processing, scoring engine, "
            "and results display update in real-time.",
            
            # Slide 12: Demo
            "Live demonstration will show analysis of: "
            "1) 'Official government statement' (Real) "
            "2) 'Alien invasion happening now' (Fake)",
            
            # Slide 13: Future Improvements
            "Next steps include: expanding keyword database, adding source verification, "
            "and implementing machine learning models.",
            
            # Slide 14: Conclusion
            "Thank you for your attention. We welcome questions about our methodology, "
            "results, or potential applications of this technology."
        ]
def add_notes_to_presentation(self):
        """Add notes to all slides with proper error handling"""
        try:
            # Check if input file exists
            if not os.path.exists(self.input_ppt):
                return False, f"Input file not found: {self.input_ppt}"
            
            # Check output directory permissions
            output_dir = os.path.dirname(self.output_ppt)
            if output_dir and not os.path.exists(output_dir):
                os.makedirs(output_dir)
            
            # Check if output file is writable
            if os.path.exists(self.output_ppt) and not os.access(self.output_ppt, os.W_OK):
                return False, f"Cannot write to output file: {self.output_ppt}"
            
            prs = Presentation(self.input_ppt)
            
            # Verify slide count
            if len(prs.slides) != 14:
                return False, f"Expected 14 slides, found {len(prs.slides)}"
            
            # Add notes to each slide
            for i, slide in enumerate(prs.slides):
                if i < len(self.slide_notes):
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = self.slide_notes[i]
            
            # Save the presentation
            prs.save(self.output_ppt)
            return True, f"Successfully added notes to all 14 slides\nSaved as: {self.output_ppt}"
            
        except PermissionError as e:
            return False, f"Permission denied: {str(e)}"
        except Exception as e:
            return False, f"Error: {str(e)}"

class NotesManagerApp(tk.Tk):
    """GUI application for managing PowerPoint notes"""
    
    def __init__(self):
        super().__init__()
        self.manager = PowerPointNotesManager()
        self.title("PowerPoint Notes Manager")
        self.geometry("700x500")
        self._create_widgets()
        
    def _create_widgets(self):
        """Create UI components"""
        main_frame = ttk.Frame(self, padding="20")
        main_frame.pack(expand=True, fill=tk.BOTH)
        
        # Header
        ttk.Label(
            main_frame,
            text="AI Fake News Detector\nPresentation Notes Manager",
            font=("Arial", 16, "bold"),
            justify="center"
        ).pack(pady=10)
        
        # File info
        info_frame = ttk.LabelFrame(main_frame, text=" File Information ", padding=10)
        info_frame.pack(fill=tk.X, pady=10)
        
        ttk.Label(info_frame, text=f"Input: {self.manager.input_ppt}").pack(anchor="w")
        ttk.Label(info_frame, text=f"Output: {self.manager.output_ppt}").pack(anchor="w")
        ttk.Label(info_frame, text="Total Slides: 14").pack(anchor="w")
        
        # Action buttons
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(pady=20)
        
        ttk.Button(
            btn_frame,
            text="Add Notes to Presentation",
            command=self._add_notes,
            width=25
        ).pack()
        
        # Status console
        self.console = tk.Text(
            main_frame,
            height=10,
            wrap=tk.WORD,
            state="disabled",
            font=("Consolas", 10),
            background="#f0f0f0"
        )
        self.console.pack(fill=tk.BOTH, expand=True)
        
        # Help text
        ttk.Label(
            main_frame,
            text="Note: Close PowerPoint before running this tool",
            foreground="red"
        ).pack(pady=5)
    
    def _log(self, message, is_error=False):
        """Add message to console"""
        self.console.config(state="normal")
        self.console.insert(
            tk.END, 
            f"{message}\n", 
            "error" if is_error else "normal"
        )
        self.console.tag_config("error", foreground="red")
        self.console.see(tk.END)
        self.console.config(state="disabled")
    
    def _add_notes(self):
        """Handle the add notes operation"""
        confirm = messagebox.askyesno(
            "Confirm",
            "This will add notes to all 14 slides. Continue?"
        )
        if confirm:
            success, message = self.manager.add_notes_to_presentation()
            self._log(message, not success)
            if success:
                messagebox.showinfo("Success", "Notes added successfully!")
                try:
                    os.startfile(self.manager.output_ppt)
                except Exception as e:
                    self._log(f"Could not open file: {str(e)}", True)
            else:
                messagebox.showerror("Error", message)

if __name__ == "__main__":
    app = NotesManagerApp()
    app.mainloop()